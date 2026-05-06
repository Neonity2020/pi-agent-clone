"""
Agent Harness Engineering 课件生成脚本
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# 创建演示文稿
prs = Presentation()
prs.slide_width = Cm(25.4)
prs.slide_height = Cm(14.29)

def add_title_slide(prs, title, subtitle=""):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景色块
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(0), 
        Inches(10), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    shape.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0xa8, 0xa8, 0xdc)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_list):
    """添加内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 顶部装饰条
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.8)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    shape.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    
    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.5), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "● " + item
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(0x2e, 0x2e, 0x4a)
        p.space_before = Pt(15)
        p.space_after = Pt(8)
    
    return slide

def add_section_slide(prs, section_title, section_number=""):
    """添加章节分隔页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 大背景色块
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    shape.line.fill.background()
    
    # 章节编号
    if section_number:
        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = section_number
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x6c, 0x5c, 0xff)
        p.alignment = PP_ALIGN.CENTER
    
    # 章节标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """添加两栏内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 顶部装饰条
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.8)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    shape.line.fill.background()
    
    # 页面标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    
    # 左栏标题
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(4.2), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x6c, 0x5c, 0xff)
    
    # 左栏内容
    left_content = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(4.2), Inches(5))
    tf = left_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "✓ " + item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x2e, 0x2e, 0x4a)
        p.space_before = Pt(10)
    
    # 右栏标题
    right_title_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.1), Inches(4.2), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0x6b, 0x6b)
    
    # 右栏内容
    right_content = slide.shapes.add_textbox(Inches(5.3), Inches(1.7), Inches(4.2), Inches(5))
    tf = right_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "✓ " + item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x2e, 0x2e, 0x4a)
        p.space_before = Pt(10)
    
    return slide

# ==================== 幻灯片内容 ====================

# 第1页：标题页
add_title_slide(prs, "Agent Harness Engineering", "智能体测试工程实践")

# 第2页：目录
add_content_slide(prs, "目 录", [
    "什么是 Agent",
    "Harness 概念解析",
    "Agent Harness Engineering 概述",
    "核心组件与架构",
    "测试框架设计",
    "评估指标体系",
    "安全性与可靠性测试",
    "工程实践案例",
    "总结与展望"
])

# 第3页：章节分隔 - 什么是Agent
add_section_slide(prs, "什么是 Agent", "01")

# 第4页：Agent概念
add_content_slide(prs, "Agent 的定义", [
    "Agent（智能体）是能够感知环境、做出决策并执行动作的自主实体",
    "核心能力：感知 → 推理 → 行动 → 学习",
    "与大语言模型(LLM)的区别：Agent具备规划、工具使用、记忆能力",
    "常见架构：ReAct、Plan-and-Execute、AutoGPT等",
    "应用场景：自动化任务、对话系统、代码生成、数据分析"
])

# 第5页：章节分隔 - Harness概念
add_section_slide(prs, "Harness 概念解析", "02")

# 第6页：Harness概念
add_two_column_slide(prs, "Harness 概念解析",
    "传统软件测试", 
    ["Test Harness 是测试环境的抽象", "用于自动化执行测试用例", "提供fixture和mock机制", "收集测试结果和覆盖率"],
    "Agent Harness",
    ["评估Agent能力的标准化框架", "模拟真实环境与交互场景", "量化Agent性能指标", "支持多维度评测"]
)

# 第7页：章节分隔 - 概述
add_section_slide(prs, "Agent Harness Engineering", "03")

# 第8页：概述
add_content_slide(prs, "Agent Harness Engineering 概述", [
    "定义：对Agent系统进行系统化测试、评估和验证的工程实践",
    "目标：确保Agent在复杂环境中的可靠性、安全性和有效性",
    "核心问题：如何全面评估Agent的规划能力、工具使用、决策质量",
    "挑战：开放性生成、长程规划、意图对齐、幻觉检测",
    "工程价值：提高Agent质量、加速迭代、支撑决策"
])

# 第9页：章节分隔 - 核心组件
add_section_slide(prs, "核心组件与架构", "04")

# 第10页：核心组件
add_content_slide(prs, "核心组件", [
    "Environment Simulator（环境模拟器）：模拟Agent运行的实际环境",
    "Task Generator（任务生成器）：自动生成多样化测试任务",
    "Agent Executor（执行器）：管理Agent生命周期和执行流程",
    "Metric Evaluator（评估器）：计算多维度性能指标",
    "Result Analyzer（分析器）：分析失败案例，生成诊断报告"
])

# 第11页：章节分隔 - 测试框架
add_section_slide(prs, "测试框架设计", "05")

# 第12页：测试框架
add_content_slide(prs, "测试框架设计原则", [
    "可重复性：相同输入必须产生一致结果",
    "隔离性：测试用例之间互不干扰",
    "可扩展性：支持新增任务类型和评估指标",
    "自动化：无需人工干预完成全流程",
    "可观测性：完整记录Agent推理过程和中间状态"
])

# 第13页：测试类型
add_two_column_slide(prs, "测试类型",
    "功能测试",
    ["任务完成率", "工具使用准确性", "响应质量评估", "对话连贯性"],
    "非功能测试",
    ["响应时延测试", "并发能力测试", "资源消耗评估", "错误恢复能力"]
)

# 第14页：章节分隔 - 评估指标
add_section_slide(prs, "评估指标体系", "06")

# 第15页：评估指标
add_content_slide(prs, "多维度评估指标", [
    "任务完成指标：成功率、步骤数、关键动作命中率",
    "质量指标：输出准确性、相关性、完整性、安全性",
    "效率指标：响应时间、token消耗、API调用次数",
    "鲁棒性指标：对抗样本表现、边界条件处理",
    "一致性指标：输出一致性、对齐人类偏好"
])

# 第16页：章节分隔 - 安全性测试
add_section_slide(prs, "安全性与可靠性测试", "07")

# 第17页：安全性测试
add_content_slide(prs, "安全性测试重点", [
    "Prompt Injection：检测恶意指令注入攻击",
    "Sensitive Data Leakage：防止敏感信息泄露",
    "Harmful Content：阻断有害内容生成",
    "Privacy Protection：个人信息保护能力",
    "Jailbreak Detection：越狱攻击防御能力"
])

# 第18页：可靠性测试
add_content_slide(prs, "可靠性测试重点", [
    "长时间运行稳定性",
    "异常输入容错能力",
    "依赖服务降级处理",
    "状态一致性保证",
    "资源泄漏检测"
])

# 第19页：章节分隔 - 工程实践
add_section_slide(prs, "工程实践案例", "08")

# 第20页：实践案例
add_content_slide(prs, "典型实践案例", [
    "AgentBench：综合评估多领域Agent能力",
    "ToolBench：工具使用能力专项评测",
    "MINT：多轮交互推理测试",
    "TrustGPT：价值观对齐评估",
    "自制Harness：根据业务场景定制测试框架"
])

# 第21页：实践建议
add_content_slide(prs, "工程实践建议", [
    "从小做起：先针对核心能力建立基准测试",
    "数据为王：积累高质量评测数据集",
    "自动化流水线：CI/CD集成自动评测",
    "持续迭代：定期更新评测标准和任务库",
    "社区协作：参与开源评测项目共建"
])

# 第22页：章节分隔 - 总结
add_section_slide(prs, "总结与展望", "09")

# 第23页：总结
add_content_slide(prs, "核心要点回顾", [
    "Agent Harness Engineering是Agent开发的关键环节",
    "标准化、可重复的评测框架是质量保障基础",
    "多维度评估指标覆盖功能、效率、安全性",
    "持续迭代的测试体系支撑Agent能力提升",
    "工程实践需要结合具体业务场景定制"
])

# 第24页：展望
add_content_slide(prs, "未来发展趋势", [
    "更接近人类认知评估的多模态评测",
    "动态对抗环境下的压力测试",
    "跨语言、跨文化的泛化能力评估",
    "Agent自我评测与反思能力研究",
    "评测效率优化与成本控制"
])

# 第25页：结束页
add_title_slide(prs, "谢谢观看", "Agent Harness Engineering 课件")

# 保存文件
output_path = "Agent_Harness_Engineering_课件.pptx"
prs.save(output_path)
print(f"✅ 课件已生成：{output_path}")
